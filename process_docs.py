"""
Document Processor - Converts business documents into searchable knowledge base chunks.
Run with: python process_docs.py
"""

import os
import sys
import json
import csv
import subprocess
import tempfile
import shutil
import hashlib
from pathlib import Path
from datetime import datetime

import yaml
import pdfplumber
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
import anthropic
import pytesseract
from pdf2image import convert_from_path


# =============================================================================
# SETUP
# =============================================================================

def load_config():
    """Load settings from config.yaml"""
    config_path = Path(__file__).parent / "config.yaml"
    with open(config_path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def setup_folders(config):
    """Create output folders if they don't exist"""
    output_path = Path(config["output_folder"])
    chunks_path = output_path / "chunks"
    converted_path = output_path / "converted"  # For legacy file conversions

    output_path.mkdir(parents=True, exist_ok=True)
    chunks_path.mkdir(exist_ok=True)
    converted_path.mkdir(exist_ok=True)

    return output_path, chunks_path, converted_path


def compute_file_hash(file_path):
    """Compute MD5 hash of file content for duplicate detection"""
    hash_md5 = hashlib.md5()
    try:
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    except Exception:
        return None


def compute_text_hash(text):
    """Compute hash of normalized text content for near-duplicate detection.

    This catches duplicates even when file formats differ (PDF vs DOCX)
    or when there are minor formatting differences.
    """
    if not text or len(text.strip()) < 100:
        return None
    # Normalize: lowercase, collapse whitespace, take first 5000 chars
    normalized = ' '.join(text.lower().split())[:5000]
    return hashlib.md5(normalized.encode('utf-8')).hexdigest()


def load_progress(output_path):
    """Load progress file to track which files have been processed"""
    progress_file = output_path / "progress.json"
    if progress_file.exists():
        with open(progress_file, "r", encoding="utf-8") as f:
            data = json.load(f)
            # Ensure required fields exist (for backward compatibility)
            if "content_hashes" not in data:
                data["content_hashes"] = {}
            if "text_hashes" not in data:
                data["text_hashes"] = {}  # For near-duplicate detection
            if "duplicates_skipped" not in data:
                data["duplicates_skipped"] = []
            return data
    return {"processed_files": [], "content_hashes": {}, "text_hashes": {}, "duplicates_skipped": [], "total_chunks": 0, "total_cost": 0.0}


def save_progress(output_path, progress):
    """Save progress to file"""
    progress_file = output_path / "progress.json"
    with open(progress_file, "w", encoding="utf-8") as f:
        json.dump(progress, f, indent=2)


# =============================================================================
# FILE DISCOVERY
# =============================================================================

def find_documents(source_folder, file_types):
    """Find all documents in the source folder and subfolders"""
    source_path = Path(source_folder)
    documents = []

    for file_type in file_types:
        for file_path in source_path.rglob(f"*{file_type}"):
            # Skip temporary files (start with ~$)
            if file_path.name.startswith("~$"):
                continue
            documents.append(file_path)

    return sorted(documents, key=lambda x: x.name)


# =============================================================================
# LEGACY FILE CONVERSION (LibreOffice)
# =============================================================================

def convert_legacy_file(file_path, converted_path, libreoffice_path):
    """Convert .doc or .ppt to .docx or .pptx using LibreOffice"""

    if not Path(libreoffice_path).exists():
        print(f"  ERROR: LibreOffice not found at {libreoffice_path}")
        print("  Please install LibreOffice or update the path in config.yaml")
        return None

    # Determine output format
    suffix = file_path.suffix.lower()
    if suffix == ".doc":
        output_filter = "MS Word 2007 XML"
        new_suffix = ".docx"
    elif suffix == ".ppt":
        output_filter = "Impress MS PowerPoint 2007 XML"
        new_suffix = ".pptx"
    else:
        return file_path  # Already modern format

    # Create output filename
    output_file = converted_path / (file_path.stem + new_suffix)

    # Skip if already converted
    if output_file.exists():
        return output_file

    # Run LibreOffice conversion
    try:
        cmd = [
            libreoffice_path,
            "--headless",
            "--convert-to", new_suffix[1:],  # Remove the dot
            "--outdir", str(converted_path),
            str(file_path)
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=60)

        if output_file.exists():
            return output_file
        else:
            print(f"  WARNING: Conversion failed for {file_path.name}")
            return None

    except subprocess.TimeoutExpired:
        print(f"  WARNING: Conversion timed out for {file_path.name}")
        return None
    except Exception as e:
        print(f"  WARNING: Conversion error for {file_path.name}: {e}")
        return None


# =============================================================================
# TEXT EXTRACTION
# =============================================================================

def extract_text_from_pdf_ocr(file_path, tesseract_path, poppler_path=None):
    """Extract text from scanned PDF using OCR"""
    try:
        # Set tesseract path
        pytesseract.pytesseract.tesseract_cmd = tesseract_path

        # Convert PDF pages to images
        print("  Using OCR (scanned PDF)...")
        if poppler_path:
            images = convert_from_path(file_path, dpi=200, poppler_path=poppler_path)
        else:
            images = convert_from_path(file_path, dpi=200)

        text_parts = []
        for i, image in enumerate(images):
            page_text = pytesseract.image_to_string(image)
            if page_text.strip():
                text_parts.append(page_text)

        return "\n\n".join(text_parts)

    except Exception as e:
        print(f"  WARNING: OCR failed: {e}")
        return ""


def extract_text_from_pdf(file_path, tesseract_path=None, poppler_path=None):
    """Extract text from PDF file, with OCR fallback for scanned PDFs"""
    text_parts = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
    except Exception as e:
        print(f"  WARNING: Could not extract text from PDF: {e}")
        return ""

    text = "\n\n".join(text_parts)

    # If no text extracted and tesseract is available, try OCR
    if not text.strip() and tesseract_path and Path(tesseract_path).exists():
        text = extract_text_from_pdf_ocr(file_path, tesseract_path, poppler_path)

    return text


def extract_text_from_docx(file_path):
    """Extract text from Word document"""
    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        print(f"  WARNING: Could not extract text from DOCX: {e}")
        return ""


def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation"""
    text_parts = []
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    except Exception as e:
        print(f"  WARNING: Could not extract text from PPTX: {e}")
        return ""
    return "\n\n".join(text_parts)


def extract_text_from_txt(file_path):
    """Extract text from plain text file"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception as e:
            print(f"  WARNING: Could not read text file: {e}")
            return ""
    except Exception as e:
        print(f"  WARNING: Could not read text file: {e}")
        return ""


def extract_text(file_path, converted_path, libreoffice_path, tesseract_path=None, poppler_path=None):
    """Extract text from any supported document type"""
    suffix = file_path.suffix.lower()

    # Handle legacy formats - convert first
    if suffix in [".doc", ".ppt"]:
        converted_file = convert_legacy_file(file_path, converted_path, libreoffice_path)
        if converted_file is None:
            return ""
        file_path = converted_file
        suffix = file_path.suffix.lower()

    # Extract based on format
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path, tesseract_path, poppler_path)
    elif suffix == ".docx":
        return extract_text_from_docx(file_path)
    elif suffix == ".pptx":
        return extract_text_from_pptx(file_path)
    elif suffix == ".txt":
        return extract_text_from_txt(file_path)
    else:
        print(f"  WARNING: Unsupported file type: {suffix}")
        return ""


# =============================================================================
# CHUNKING
# =============================================================================

def count_words(text):
    """Count words in text"""
    return len(text.split())


def create_chunks(text, min_words, max_words, overlap_words):
    """Split text into chunks of target size with overlap"""

    if not text.strip():
        return []

    # Split into paragraphs
    paragraphs = [p.strip() for p in text.split("\n") if p.strip()]

    if not paragraphs:
        return []

    chunks = []
    current_chunk = []
    current_word_count = 0

    for para in paragraphs:
        para_words = count_words(para)

        # If adding this paragraph exceeds max, save current chunk and start new one
        if current_word_count + para_words > max_words and current_chunk:
            chunk_text = "\n\n".join(current_chunk)
            chunks.append(chunk_text)

            # Keep last bit for overlap
            overlap_text = []
            overlap_count = 0
            for p in reversed(current_chunk):
                p_words = count_words(p)
                if overlap_count + p_words <= overlap_words:
                    overlap_text.insert(0, p)
                    overlap_count += p_words
                else:
                    break

            current_chunk = overlap_text
            current_word_count = overlap_count

        current_chunk.append(para)
        current_word_count += para_words

    # Don't forget the last chunk
    if current_chunk:
        chunk_text = "\n\n".join(current_chunk)
        # Only add if it has meaningful content (at least 50 words)
        if count_words(chunk_text) >= 50:
            chunks.append(chunk_text)

    return chunks


# =============================================================================
# CLAUDE API - TOPIC CLASSIFICATION
# =============================================================================

def classify_chunk(client, chunk_text, source_filename, model):
    """Use Claude to classify the chunk and generate metadata"""

    prompt = f"""Analyze this text chunk from a business document and provide:
1. A topic category (2-4 words, e.g., "Financial Planning", "Employee Training", "Sales Strategy")
2. A one-sentence summary (max 20 words)
3. 3-5 relevant keywords

Source file: {source_filename}

Text to analyze:
{chunk_text[:3000]}

Respond in this exact JSON format (no other text):
{{"topic": "...", "summary": "...", "keywords": ["...", "...", "..."]}}"""

    try:
        response = client.messages.create(
            model=model,
            max_tokens=200,
            messages=[{"role": "user", "content": prompt}]
        )

        # Parse response
        response_text = response.content[0].text.strip()

        # Calculate cost (Haiku pricing as of 2024)
        input_tokens = response.usage.input_tokens
        output_tokens = response.usage.output_tokens
        cost = (input_tokens * 0.00025 / 1000) + (output_tokens * 0.00125 / 1000)

        # Parse JSON from response
        metadata = json.loads(response_text)
        metadata["api_cost"] = cost

        return metadata

    except json.JSONDecodeError:
        # If Claude doesn't return valid JSON, use defaults
        return {
            "topic": "Uncategorized",
            "summary": "Could not generate summary",
            "keywords": [],
            "api_cost": 0.001  # Estimate
        }
    except Exception as e:
        print(f"  WARNING: API error: {e}")
        return {
            "topic": "Uncategorized",
            "summary": "API error during classification",
            "keywords": [],
            "api_cost": 0.0
        }


def detect_presentation_topic(client, text, source_filename, model):
    """Detect the overall topic/theme of a presentation.

    This is used to provide context for all chunks from the presentation,
    so that detail slides remain searchable by the main topic.
    """
    import re

    # Use first ~4000 chars which typically covers intro/overview slides
    sample_text = text[:4000]

    prompt = f"""What is the MAIN TOPIC of this presentation in 2-5 words?

Examples of good responses:
- Customer Focus
- Leadership Development
- Inventory Management
- Change Management

Presentation: {source_filename}

Text:
{sample_text}

Main topic (2-5 words only):"""

    try:
        response = client.messages.create(
            model=model,
            max_tokens=30,
            messages=[{"role": "user", "content": prompt}]
        )

        topic = response.content[0].text.strip()
        # Clean up common patterns
        topic = re.sub(r'^(The main topic is|Main topic:|Topic:)\s*', '', topic, flags=re.IGNORECASE)
        topic = topic.strip('"\'.,')

        # Calculate cost
        input_tokens = response.usage.input_tokens
        output_tokens = response.usage.output_tokens
        cost = (input_tokens * 0.00025 / 1000) + (output_tokens * 0.00125 / 1000)

        return topic, cost

    except Exception as e:
        print(f"  WARNING: Could not detect presentation topic: {e}")
        return None, 0.0


# =============================================================================
# OUTPUT
# =============================================================================

def save_chunk(chunks_path, chunk_text, chunk_id, metadata, source_file, subfolder):
    """Save chunk as text file"""
    filename = f"{chunk_id}.txt"
    file_path = chunks_path / filename

    # Create header with metadata
    header = f"""SOURCE: {source_file}
SUBFOLDER: {subfolder}
TOPIC: {metadata['topic']}
SUMMARY: {metadata['summary']}
KEYWORDS: {', '.join(metadata['keywords'])}
---

"""

    with open(file_path, "w", encoding="utf-8") as f:
        f.write(header + chunk_text)

    return filename


def init_index_csv(output_path):
    """Create the index CSV with headers"""
    index_path = output_path / "index.csv"
    if not index_path.exists():
        with open(index_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow([
                "chunk_id", "chunk_file", "source_file", "subfolder",
                "chunk_number", "word_count", "topic", "summary",
                "keywords", "created_date"
            ])
    return index_path


def append_to_index(index_path, chunk_id, chunk_file, source_file, subfolder,
                    chunk_number, word_count, metadata):
    """Add a row to the index CSV"""
    with open(index_path, "a", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow([
            chunk_id,
            chunk_file,
            source_file,
            subfolder,
            chunk_number,
            word_count,
            metadata["topic"],
            metadata["summary"],
            ", ".join(metadata["keywords"]),
            datetime.now().isoformat()
        ])


# =============================================================================
# MAIN PROCESSING
# =============================================================================

def main():
    print("\n" + "="*60)
    print("DOCUMENT PROCESSOR")
    print("="*60 + "\n")

    # Load environment and config
    load_dotenv()
    api_key = os.getenv("ANTHROPIC_API_KEY")

    if not api_key or api_key == "sk-ant-your-key-here":
        print("ERROR: Please add your Anthropic API key to the .env file")
        print("1. Copy .env.example to .env")
        print("2. Replace 'sk-ant-your-key-here' with your actual API key")
        sys.exit(1)

    config = load_config()

    # Setup
    output_path, chunks_path, converted_path = setup_folders(config)
    progress = load_progress(output_path)
    index_path = init_index_csv(output_path)

    # Initialize Claude client
    client = anthropic.Anthropic(api_key=api_key)

    # Find documents
    documents = find_documents(config["source_folder"], config["file_types"])

    # Apply test mode limit
    if config.get("test_mode"):
        documents = documents[:config["test_mode"]]
        print(f"TEST MODE: Processing only {config['test_mode']} files\n")

    # Filter out already processed
    remaining = [d for d in documents if str(d) not in progress["processed_files"]]

    print(f"Found {len(documents)} documents total")
    print(f"Already processed: {len(progress['processed_files'])}")
    if progress["duplicates_skipped"]:
        print(f"  (includes {len(progress['duplicates_skipped'])} duplicates skipped)")
    print(f"Remaining to process: {len(remaining)}")
    print(f"Output folder: {output_path}\n")

    if not remaining:
        print("All documents have been processed!")
        return

    # Get source folder for calculating relative paths
    source_folder = Path(config["source_folder"])

    # Track duplicates found in this run
    duplicates_found = 0

    # Process each document
    for doc_num, doc_path in enumerate(remaining, 1):
        print(f"[{doc_num}/{len(remaining)}] Processing: {doc_path.name}")

        # Check for duplicate content
        file_hash = compute_file_hash(doc_path)
        if file_hash and file_hash in progress["content_hashes"]:
            original_file = progress["content_hashes"][file_hash]
            print(f"  Skipped: Duplicate content (same as {Path(original_file).name})")
            progress["processed_files"].append(str(doc_path))
            progress["duplicates_skipped"].append({
                "file": str(doc_path),
                "duplicate_of": original_file,
                "hash": file_hash
            })
            duplicates_found += 1
            save_progress(output_path, progress)
            continue

        # Get subfolder (relative to source)
        try:
            relative_path = doc_path.relative_to(source_folder)
            subfolder = str(relative_path.parent) if relative_path.parent != Path(".") else "root"
        except ValueError:
            subfolder = "root"

        # Extract text
        text = extract_text(doc_path, converted_path, config["libreoffice_path"], config.get("tesseract_path"), config.get("poppler_path"))

        if not text.strip():
            print("  Skipped: No text extracted")
            progress["processed_files"].append(str(doc_path))
            save_progress(output_path, progress)
            continue

        # Check for near-duplicate content (same text, different file format)
        text_hash = compute_text_hash(text)
        if text_hash and text_hash in progress["text_hashes"]:
            original_file = progress["text_hashes"][text_hash]
            print(f"  Skipped: Near-duplicate content (same text as {Path(original_file).name})")
            progress["processed_files"].append(str(doc_path))
            progress["duplicates_skipped"].append({
                "file": str(doc_path),
                "duplicate_of": original_file,
                "type": "near-duplicate",
                "hash": text_hash
            })
            save_progress(output_path, progress)
            continue

        # Create chunks
        chunks = create_chunks(
            text,
            config["chunk_min_words"],
            config["chunk_max_words"],
            config["chunk_overlap_words"]
        )

        if not chunks:
            print("  Skipped: No chunks created (document too short)")
            progress["processed_files"].append(str(doc_path))
            save_progress(output_path, progress)
            continue

        print(f"  Created {len(chunks)} chunks")

        # For PowerPoint files, detect the overall presentation topic
        # This ensures detail slides remain searchable by the main topic
        presentation_topic = None
        doc_cost = 0.0
        if doc_path.suffix.lower() in [".pptx", ".ppt"]:
            presentation_topic, topic_cost = detect_presentation_topic(
                client, text, doc_path.name, config["claude_model"]
            )
            doc_cost += topic_cost
            if presentation_topic:
                print(f"  Presentation topic: {presentation_topic}")

        # Process each chunk
        for chunk_num, chunk_text in enumerate(chunks, 1):
            chunk_id = f"{doc_path.stem}_chunk{chunk_num:03d}"

            # Classify with Claude
            metadata = classify_chunk(
                client,
                chunk_text,
                doc_path.name,
                config["claude_model"]
            )
            doc_cost += metadata.get("api_cost", 0)

            # Add contextual header to chunk content for better searchability
            # This embeds the document name and presentation topic in the searchable text
            if presentation_topic:
                context_header = f"[Document: {doc_path.name}]\n[Presentation Topic: {presentation_topic}]\n\n"
            else:
                context_header = f"[Document: {doc_path.name}]\n\n"
            chunk_text_with_context = context_header + chunk_text

            # Save chunk file
            chunk_file = save_chunk(
                chunks_path, chunk_text_with_context, chunk_id, metadata,
                doc_path.name, subfolder
            )

            # Update index
            append_to_index(
                index_path, chunk_id, chunk_file, doc_path.name,
                subfolder, chunk_num, count_words(chunk_text), metadata
            )

            progress["total_chunks"] += 1

        # Update progress and store content hashes
        progress["processed_files"].append(str(doc_path))
        if file_hash:
            progress["content_hashes"][file_hash] = str(doc_path)
        if text_hash:
            progress["text_hashes"][text_hash] = str(doc_path)
        progress["total_cost"] += doc_cost
        save_progress(output_path, progress)

        print(f"  Cost for this doc: ${doc_cost:.4f}")

    # Final summary
    print("\n" + "="*60)
    print("PROCESSING COMPLETE")
    print("="*60)
    print(f"Total documents processed: {len(progress['processed_files'])}")
    print(f"Total chunks created: {progress['total_chunks']}")
    print(f"Total API cost: ${progress['total_cost']:.4f}")
    if duplicates_found > 0:
        print(f"Duplicates skipped (this run): {duplicates_found}")
    if progress["duplicates_skipped"]:
        print(f"Total duplicates skipped (all time): {len(progress['duplicates_skipped'])}")
    print(f"\nOutput location: {output_path}")
    print(f"  - Chunks: {chunks_path}")
    print(f"  - Index: {index_path}")


if __name__ == "__main__":
    main()
